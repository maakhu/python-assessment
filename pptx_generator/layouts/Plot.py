import csv
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
import matplotlib.pyplot as plt


class PlotSlide:
    def __init__(self, title, content, config):
        self.title = title
        self.content = content
        self.config = config

    def create_slide(self, presentation):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        title = slide.shapes.title
        title.text = self.title
        # picture_placeholder = slide.placeholders[0]

        left = Inches(1)
        top = Inches(1.5)
        width = Inches(6)
        height = Inches(4.5)

        chart_data = CategoryChartData()
        with open(self.content, 'r') as file:
            for line in file:
                x, y = line.strip().split(';')
                chart_data.append((f"X: {x}", float(y)))

        # Add the chart to the slide
        x, y, cx, cy = Inches(1), Inches(3), Inches(8), Inches(4.5)
        slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
        )
        #
        # x_values = self.config['x-label']
        # y_values = self.config['x-label']
        #
        # # Generate the plot
        # x_data, y_data = self.read_plot_data(self.content)
        # plt.plot(x_data, y_data)
        # plt.xlabel(x_values)
        # plt.ylabel(y_values)
        # plt.title(self.title)
        #
        # # Save the plot as an image
        # plot_filename = "plot.png"
        # plt.savefig(plot_filename)
        # plt.close()

        # Insert the plot image into the slide
        # picture = picture_placeholder.insert_picture(plot_filename)
        # picture.left = left
        # picture.top = top
        # picture.width = width
        # picture.height = height

    @staticmethod
    def read_plot_data(filename):
        x_data = []
        y_data = []
        with open(filename, 'r') as file:
            reader = csv.reader(file, delimiter=';')
            for row in reader:
                x_data.append(float(row[0]))
                y_data.append(float(row[1]))
        return x_data, y_data
