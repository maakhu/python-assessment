from pptx.util import Inches, Pt


class TextSlide:
    def __init__(self, title, content):
        self.title = title
        self.content = content

    def create_slide(self, presentation):
        slide_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(slide_layout)

        left = top = width = height = Inches(1)
        tx_box = slide.shapes.add_textbox(left, top, width, height)
        tf = tx_box.text_frame

        p = tf.add_paragraph()
        p.text = self.title
        p.font.size = Pt(40)

        p = tf.add_paragraph()
        p.text = self.content

