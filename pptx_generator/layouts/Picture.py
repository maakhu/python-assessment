from pptx.util import Inches

class PictureSlide:
    def __init__(self, title, img):
        self.title = title
        self.img = img

    def create_slide(self, presentation):
        img_path = self.img
        slide_layout = presentation.slide_layouts[5]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = self.title

        top = Inches(1.5)
        left = Inches(1)
        height = Inches(5.5)
        pic = slide.shapes.add_picture(img_path, left, top, height=height)

