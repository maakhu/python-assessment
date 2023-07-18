import unittest
from pptx import Presentation
from pptx_generator import TitleSlide


class TestTitleSlide(unittest.TestCase):
    def test_title_slide_creation(self):
        # Create a new presentation
        presentation = Presentation()

        # Create a TitleSlide object
        title_slide = TitleSlide("Test Title", "Test Subtitle")

        # Call the create_slide method to generate the slide
        title_slide.create_slide(presentation)

        # Get the generated slide
        slide = presentation.slides[0]

        # Assertions to check if the title and subtitle are correctly set on the slide
        self.assertEqual(slide.shapes.title.text, "Test Title")
        self.assertEqual(slide.placeholders[1].text, "Test Subtitle")

        # Additional assertions for other properties, if necessary
        # For example, if you want to check the slide layout used:
        title_slide_layout = presentation.slide_layouts[0]
        self.assertEqual(slide.slide_layout, title_slide_layout)

    def test_title_slide_font_size(self):
        # Create a new presentation
        presentation = Presentation()

        # Create a TitleSlide object
        title_slide = TitleSlide("Test Title", "Test Subtitle")

        # Call the create_slide method to generate the slide
        title_slide.create_slide(presentation)

        # Get the generated slide
        slide = presentation.slides[0]
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        # Assertions to check the font size of the title and subtitle
        expected_title_font_size = 40  # Assuming the default font size is 40 Pt
        expected_subtitle_font_size = 24  # Assuming the default font size is 24 Pt

        self.assertEqual(title_shape.text_frame.paragraphs[0].runs[0].font.size, expected_title_font_size)
        self.assertEqual(subtitle_shape.text_frame.paragraphs[0].runs[0].font.size, expected_subtitle_font_size)


if __name__ == '__main__':
    unittest.main()
