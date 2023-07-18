import unittest
from pptx import Presentation
from pptx_generator import TitleSlide


class TestTitleSlide(unittest.TestCase):
    def test_title_slide_creation(self):
        presentation = Presentation()
        title_slide = TitleSlide("Test Title", "Test Subtitle")
        title_slide.create_slide(presentation)

        slide = presentation.slides[0]

        self.assertEqual(slide.shapes.title.text, "Test Title")
        self.assertEqual(slide.placeholders[1].text, "Test Subtitle")

        title_slide_layout = presentation.slide_layouts[0]
        self.assertEqual(slide.slide_layout, title_slide_layout)


if __name__ == '__main__':
    unittest.main()
